from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

import fitz
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
}


def clean_text(text: str) -> str:
    """Remove unnecessary spaces and blank lines."""

    cleaned_lines = []

    for line in text.splitlines():
        cleaned_line = " ".join(line.split())

        if cleaned_line:
            cleaned_lines.append(cleaned_line)

    return "\n".join(cleaned_lines)


def split_text(
    text: str,
    chunk_size: int = 2200,
    overlap: int = 250,
) -> list[str]:
    """Split text into overlapping chunks."""

    if chunk_size <= overlap:
        raise ValueError(
            "chunk_size must be greater than overlap."
        )

    text = text.strip()

    if not text:
        return []

    chunks = []
    start = 0

    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end].strip()

        if chunk:
            chunks.append(chunk)

        if end == len(text):
            break

        start += chunk_size - overlap

    return chunks


def create_records(
    document_name: str,
    locator_type: str,
    locator_number: int,
    text: str,
) -> list[dict[str, Any]]:
    """Turn one page, slide or paragraph into searchable records."""

    cleaned_text = clean_text(text)

    if not cleaned_text:
        return []

    records = []

    for chunk_number, chunk_text in enumerate(
        split_text(cleaned_text),
        start=1,
    ):
        records.append(
            {
                "document": document_name,
                "locator_type": locator_type,
                "locator_number": locator_number,
                "chunk": chunk_number,
                "text": chunk_text,
            }
        )

    return records


def process_pdf(
    file_bytes: bytes,
    document_name: str,
) -> list[dict[str, Any]]:
    """Extract text from PDF pages."""

    records = []

    pdf = fitz.open(
        stream=file_bytes,
        filetype="pdf",
    )

    try:
        for page_index in range(pdf.page_count):
            page = pdf.load_page(page_index)

            page_text = page.get_text(
                "text",
                sort=True,
            )

            records.extend(
                create_records(
                    document_name=document_name,
                    locator_type="page",
                    locator_number=page_index + 1,
                    text=page_text,
                )
            )

    finally:
        pdf.close()

    return records


def process_word(
    file_bytes: bytes,
    document_name: str,
) -> list[dict[str, Any]]:
    """Extract text from Word paragraphs and tables."""

    records = []
    word_document = Document(BytesIO(file_bytes))

    for paragraph_number, paragraph in enumerate(
        word_document.paragraphs,
        start=1,
    ):
        if paragraph.text.strip():
            records.extend(
                create_records(
                    document_name=document_name,
                    locator_type="paragraph",
                    locator_number=paragraph_number,
                    text=paragraph.text,
                )
            )

    table_number = 0

    for table in word_document.tables:
        table_number += 1

        table_rows = []

        for row in table.rows:
            row_values = [
                cell.text.strip()
                for cell in row.cells
            ]

            table_rows.append(
                " | ".join(row_values)
            )

        table_text = "\n".join(table_rows)

        records.extend(
            create_records(
                document_name=document_name,
                locator_type="table",
                locator_number=table_number,
                text=table_text,
            )
        )

    return records


def process_powerpoint(
    file_bytes: bytes,
    document_name: str,
) -> list[dict[str, Any]]:
    """Extract text from PowerPoint slides."""

    records = []
    presentation = Presentation(BytesIO(file_bytes))

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    row_values = [
                        cell.text.strip()
                        for cell in row.cells
                    ]

                    slide_text.append(
                        " | ".join(row_values)
                    )

        records.extend(
            create_records(
                document_name=document_name,
                locator_type="slide",
                locator_number=slide_number,
                text="\n".join(slide_text),
            )
        )

    return records


def process_text_file(
    file_bytes: bytes,
    document_name: str,
) -> list[dict[str, Any]]:
    """Extract text from TXT or Markdown files."""

    text = file_bytes.decode(
        "utf-8",
        errors="replace",
    )

    records = []

    for section_number, chunk in enumerate(
        split_text(
            clean_text(text),
            chunk_size=2200,
            overlap=250,
        ),
        start=1,
    ):
        records.append(
            {
                "document": document_name,
                "locator_type": "section",
                "locator_number": section_number,
                "chunk": 1,
                "text": chunk,
            }
        )

    return records


def process_document(
    file_bytes: bytes,
    document_name: str,
) -> list[dict[str, Any]]:
    """Choose the correct document processor."""

    extension = Path(document_name).suffix.lower()

    if extension == ".pdf":
        return process_pdf(
            file_bytes,
            document_name,
        )

    if extension == ".docx":
        return process_word(
            file_bytes,
            document_name,
        )

    if extension == ".pptx":
        return process_powerpoint(
            file_bytes,
            document_name,
        )

    if extension in {".txt", ".md"}:
        return process_text_file(
            file_bytes,
            document_name,
        )

    raise ValueError(
        f"Unsupported file type: {extension}"
    )